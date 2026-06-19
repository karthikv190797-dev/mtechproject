from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
COLOR_DARK_BLUE = RGBColor(25, 50, 100)
COLOR_LIGHT_BLUE = RGBColor(70, 130, 180)
COLOR_ACCENT = RGBColor(220, 20, 60)
COLOR_TEXT = RGBColor(40, 40, 40)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DARK_BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(54)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = RGBColor(255, 255, 255)
    subtitle_p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_list):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_DARK_BLUE
    
    # Add blue line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9), Inches(0))
    line.line.color.rgb = COLOR_LIGHT_BLUE
    line.line.width = Pt(3)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)

# Slide 1: Title Slide
add_title_slide(prs, "Self-Correcting Financial\nIntelligence System", 
                "End-to-End Multi-Agent RAG Pipeline\nM.Tech Dissertation\nKarthik V | May 2026")

# Slide 2: Project Overview
add_content_slide(prs, "Project Overview", [
    "• Dissertation: Self-Correcting Financial Intelligence Systems",
    "• Focus: Natural Language Processing, Information Retrieval, Financial Document Understanding",
    "• Organization: Voya Global Services Pvt Ltd",
    "• Supervisor: ShivaChandran Masilamani (Head of Investments)",
    "• Program: M.Tech in Artificial Intelligence and Machine Learning",
    "• Key Goal: Trustworthy, explainable AI for enterprise financial analytics"
])

# Slide 3: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "• Financial analysts manually query fragmented data ecosystems (Snowflake, SEC filings, audit docs)",
    "• Current RAG systems generate hallucinated/unverifiable results",
    "• Lack domain-adapted embeddings for financial entities and numerical tables",
    "• Privacy concerns with cloud-hosted LLMs in enterprise contexts",
    "• Static single-retriever pipelines cannot identify or correct their own errors",
    "• Manual validation creates latency, redundancy, and inconsistent results"
])

# Slide 4: Objectives
add_content_slide(prs, "Core Objectives", [
    "✓ Design multi-agent RAG framework with LangGraph orchestration",
    "✓ Implement Teacher-Student distillation for optimized tabular reasoning",
    "✓ Develop hybrid retrieval (FAISS + BM25) for structured & unstructured data",
    "✓ Build Evaluator agent for intelligent self-correction",
    "✓ Design Cross-Model Verification engine for consensus validation",
    "✓ Deploy as containerized, privacy-compliant on-premise solution"
])

# Slide 5: System Architecture
add_content_slide(prs, "System Architecture Components", [
    "1. Retrieval Layer: Hybrid search (Semantic FAISS + Lexical BM25)",
    "2. Reasoning Layer: Teacher-Student distillation with fine-tuned embeddings",
    "3. Verification Layer: Multi-model cross-verification for hallucination detection",
    "4. Evaluation Layer: Evaluator agent assessing faithfulness & numerical consistency",
    "5. Orchestration Layer: LangGraph-based workflow engine",
    "6. Backend: FastAPI REST APIs for scalable async processing",
    "7. Frontend: Streamlit UI for explainability and user trust"
])

# Slide 6: End-to-End Pipeline
add_content_slide(prs, "End-to-End Pipeline", [
    "STEP 1: Data Ingestion → SEC 10K + Snowflake databases",
    "STEP 2: Hybrid Retrieval → FAISS (semantic) + BM25 (keyword) search",
    "STEP 3: Reader Agent → Extract context and prepare evidence",
    "STEP 4: Teacher Model → Deep understanding of tables and text",
    "STEP 5: Student Model → Lightweight, optimized reasoning",
    "STEP 6: Evaluator Assessment → Confidence scoring and consistency checks",
    "STEP 7: Cross-Verification → Multiple LLMs validate consensus",
    "STEP 8: Self-Correction Loop → Re-retrieve if confidence < threshold"
])

# Slide 7: Key Working Models
add_content_slide(prs, "Key Working Models", [
    "• Retrieval-Augmented Generation (RAG): Combines context retrieval with LLM reasoning",
    "• Multi-Agent Orchestration: Retriever, Reader, Verifier, Evaluator agents collaborate",
    "• Teacher-Student Distillation: Large model supervises compact model training",
    "• Evaluator-in-the-Loop: Continuous assessment with automatic error correction",
    "• Cross-Model Verification: Consensus voting across multiple LLMs",
    "• Hybrid Retrieval Scoring: Weighted combination of semantic + lexical relevance"
])

# Slide 8: Technical Implementation
add_content_slide(prs, "Technical Implementation", [
    "🔹 Data Ingestion: PDFium2 for SEC filings, native Snowflake connector",
    "🔹 Vector Search: FAISS for efficient semantic similarity at scale",
    "🔹 Keyword Matching: BM25 probabilistic relevance framework",
    "🔹 Embedding Models: Domain-adapted embeddings from SEC financial datasets",
    "🔹 LLM Integration: Open-source models (Ollama, GLM-130B) for on-premise deployment",
    "🔹 Orchestration: LangGraph for dynamic task routing and feedback propagation"
])

# Slide 9: Hybrid Retrieval Strategy
add_content_slide(prs, "Hybrid Retrieval: FAISS + BM25", [
    "FAISS (Dense Vector Search):",
    "  • Semantic understanding via embeddings",
    "  • Captures meaning and context of financial queries",
    "",
    "BM25 (Sparse Keyword Matching):",
    "  • Exact keyword matching for financial terminology",
    "  • Ensures lexical precision for specific entities",
    "",
    "Weighted Scoring: Combined score = α·FAISS_score + β·BM25_score"
])

# Slide 10: Evaluator Self-Correction Loop
add_content_slide(prs, "Evaluator Self-Correction Mechanism", [
    "Multi-Criteria Scoring:",
    "  ✓ Relevance Score: Is answer addressing the query?",
    "  ✓ Factual Consistency: Can results be verified against evidence?",
    "  ✓ Numerical Coherence: Are financial figures mathematically sound?",
    "  ✓ Entailment Score: Does evidence logically support the answer?",
    "",
    "Auto-Correction Triggers:",
    "  → If confidence < threshold → Query reformulation",
    "  → Trigger re-retrieval or passage reranking"
])

# Slide 11: Deployment & Infrastructure
add_content_slide(prs, "Deployment & Infrastructure", [
    "🐳 Containerization: Docker for reproducible, portable deployments",
    "⚙️ Configuration: YAML-based orchestrator for easy customization",
    "🔐 Data Privacy: On-premise deployment, no external data transmission",
    "📡 Backend: FastAPI for RESTful APIs with async processing",
    "🎨 Frontend: Streamlit for interactive queries & confidence visualization",
    "📊 Database: Native Snowflake integration with structured/unstructured data"
])

# Slide 12: Evaluation Metrics
add_content_slide(prs, "Evaluation Metrics", [
    "Retrieval Performance:",
    "  • Precision@K, Mean Reciprocal Rank (MRR), NDCG",
    "",
    "Reasoning Performance:",
    "  • Faithfulness Score, Hallucination Detection Accuracy",
    "  • Numerical Consistency, Semantic Relevance",
    "",
    "System-level Performance:",
    "  • Cross-Model Agreement Rate",
    "  • End-to-end accuracy on financial question answering tasks",
    "  • Inference latency and resource utilization"
])

# Slide 13: Project Timeline
add_content_slide(prs, "Project Timeline", [
    "Phase 1 (Apr 25 - May 10): Literature review & abstract submission",
    "Phase 2 (May 11 - Jun 4): Data ingestion & baseline RAG implementation",
    "Phase 3 (Jun 5 - Jun 21): Multi-agent orchestration & teacher-student distillation",
    "Phase 4 (Jun 22 - Jul 15): Verifier & evaluator feedback loop integration",
    "Phase 5 (Jul 16 - Aug 2): Docker packaging & benchmark evaluations",
    "Phase 6 (Aug 3 - Aug 28): Final presentation & VIVA examination"
])

# Slide 14: Key Deliverables
add_content_slide(prs, "Key Deliverables", [
    "✅ Multi-agent orchestration framework (LangGraph)",
    "✅ Hybrid retrieval pipeline (FAISS + BM25 integration)",
    "✅ Teacher-Student distillation module for table reasoning",
    "✅ Evaluator-driven self-correction mechanism",
    "✅ Cross-model verification engine",
    "✅ FastAPI backend with REST endpoints",
    "✅ Streamlit interactive UI with explainability",
    "✅ Docker containerized deployment package"
])

# Slide 15: Expected Outcomes
add_content_slide(prs, "Expected Outcomes & Impact", [
    "🎯 Trustworthy Financial Intelligence: Verifiable, explainable AI for financial analytics",
    "🎯 Enterprise Deployment: Private, compliant, on-premise solution",
    "🎯 Self-Improving System: Automatic error detection and correction capability",
    "🎯 Scalable Architecture: Modular design supporting multiple use cases",
    "🎯 Publication Potential: Research contributions in RAG, multi-agent systems, evaluator frameworks",
    "🎯 Industry Adoption: Ready for deployment in financial services enterprises"
])

# Slide 16: Conclusion
add_content_slide(prs, "Conclusion", [
    "Self-Correcting Financial Intelligence System combines:",
    "  • Multi-agent retrieval and reasoning",
    "  • Teacher-student knowledge transfer",
    "  • Continuous evaluator-driven feedback",
    "  • Cross-model verification for trustworthiness",
    "",
    "Result: An intelligent, self-improving financial assistant",
    "Suitable for enterprise deployment with explainability and regulatory compliance"
])

# Save presentation
prs.save('Self-Correcting_Financial_Intelligence_System.pptx')
print("✅ PowerPoint presentation created successfully!")
print("📄 File: Self-Correcting_Financial_Intelligence_System.pptx")
print("📊 Total slides: 16")
