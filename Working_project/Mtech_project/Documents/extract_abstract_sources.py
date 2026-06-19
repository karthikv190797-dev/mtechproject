from pathlib import Path
import json
from pypdf import PdfReader
from pptx import Presentation


def extract_pdf_text(pdf_path: Path) -> str:
    reader = PdfReader(str(pdf_path))
    chunks = []
    for i, page in enumerate(reader.pages, 1):
        text = page.extract_text() or ""
        chunks.append(f"\n--- PAGE {i} ---\n{text}")
    return "\n".join(chunks)


def extract_pptx_text(pptx_path: Path) -> str:
    prs = Presentation(str(pptx_path))
    slides_out = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text.strip())
        body = "\n".join([t for t in texts if t])
        slides_out.append(f"\n--- SLIDE {i} ---\n{body}")
    return "\n".join(slides_out)


def main() -> None:
    root = Path(r"e:\Final_Mtech_Project\Mtech_project")
    docs = root / "Documents"

    pdf_path = docs / "Final_Project abstract document.pdf"
    pptx_path = docs / "Self-Correcting_Financial_Intelligence_System.pptx"

    output_dir = root / "Documents" / "extracted"
    output_dir.mkdir(parents=True, exist_ok=True)

    out = {}

    if pdf_path.exists():
        out["pdf_text"] = extract_pdf_text(pdf_path)
    else:
        out["pdf_text"] = "PDF file not found."

    if pptx_path.exists():
        out["pptx_text"] = extract_pptx_text(pptx_path)
    else:
        out["pptx_text"] = "PPTX file not found."

    json_path = output_dir / "abstract_sources_extracted.json"
    txt_path = output_dir / "abstract_sources_extracted.txt"

    json_path.write_text(json.dumps(out, indent=2, ensure_ascii=False), encoding="utf-8")

    flat = []
    flat.append("=== PDF EXTRACT ===")
    flat.append(out["pdf_text"])
    flat.append("\n=== PPTX EXTRACT ===")
    flat.append(out["pptx_text"])
    txt_path.write_text("\n".join(flat), encoding="utf-8")

    print(f"Saved: {json_path}")
    print(f"Saved: {txt_path}")


if __name__ == "__main__":
    main()
